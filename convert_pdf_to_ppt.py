#!/usr/bin/env python3
# -*- coding: utf-8 -*-

"""
将单依纯《歌手》节目数据分析报告PDF转换为PPT格式，并添加解说文本备注
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import pandas as pd
import glob

def create_presentation():
    """创建PPT演示文稿"""
    prs = Presentation()
    prs.slide_width = Inches(13.33)  # 16:9 宽屏比例
    prs.slide_height = Inches(7.5)
    return prs

def set_slide_background(slide, color_hex="#F0F8FF"):
    """设置幻灯片背景颜色"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor.from_string(color_hex[1:])  # 移除#号

def add_title_slide(prs):
    """添加标题幻灯片"""
    slide_layout = prs.slide_layouts[0]  # 标题幻灯片布局
    slide = prs.slides.add_slide(slide_layout)
    
    # 设置背景
    set_slide_background(slide, "#1a1a2e")
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "单依纯《歌手》节目数据分析报告"
    subtitle.text = "基于B站平台数据的专业分析"
    
    # 设置标题样式
    title_format = title.text_frame.paragraphs[0]
    title_format.font.size = Pt(32)
    title_format.font.bold = True
    title_format.font.color.rgb = RGBColor(255, 255, 255)  # 白色
    
    # 设置副标题样式
    subtitle_format = subtitle.text_frame.paragraphs[0]
    subtitle_format.font.size = Pt(24)
    subtitle_format.font.color.rgb = RGBColor(200, 200, 200)  # 浅灰色
    
    # 添加装饰性元素
    left = Inches(10)
    top = Inches(0.5)
    width = Inches(2)
    height = Inches(2)
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, width, height)
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(220, 20, 60)  # 猩红色
    shape.line.color.rgb = RGBColor(255, 255, 255)  # 白色边框
    
    # 添加备注
    slide.notes_slide.notes_text_frame.text = (
        "欢迎观看单依纯《歌手》节目数据分析报告演示。\n\n"
        "本报告基于单依纯在2025年《歌手》节目中13首现场演唱视频的数据进行分析，\n"
        "通过热度趋势、弹幕情感、关键词分析和综合评分等维度，全面评估其在节目中的表现。\n\n"
        "接下来我将逐一为您介绍各项分析结果。"
    )

def add_outline_slide(prs):
    """添加报告大纲幻灯片"""
    slide_layout = prs.slide_layouts[1]  # 标题和内容布局
    slide = prs.slides.add_slide(slide_layout)
    
    # 设置背景
    set_slide_background(slide, "#16213e")
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "报告内容大纲"
    
    # 设置标题样式
    title_format = title.text_frame.paragraphs[0]
    title_format.font.color.rgb = RGBColor(255, 255, 255)  # 白色
    
    # 添加内容
    content.text = (
        "1. 项目概述与数据概况\n"
        "2. 热度趋势分析\n"
        "3. 弹幕情感分析\n"
        "4. 关键词分析\n"
        "5. 综合表现分析\n"
        "6. 专业评价\n"
        "7. 发展建议\n"
        "8. 观众参与度分析"
    )
    
    # 设置内容样式
    for i, paragraph in enumerate(content.text_frame.paragraphs):
        paragraph.font.size = Pt(20)
        paragraph.font.color.rgb = RGBColor(220, 220, 220)  # 浅灰色
        if i == 0:  # 第一段加粗
            paragraph.font.bold = True
    
    # 添加备注
    slide.notes_slide.notes_text_frame.text = (
        "本次报告主要包含以下八个部分：\n\n"
        "首先是项目概述和数据概况，介绍分析背景和基础数据；\n"
        "然后是热度趋势分析，评估作品的受欢迎程度；\n"
        "接着是弹幕情感分析，了解观众的态度倾向；\n"
        "关键词分析帮助我们发现观众关注的焦点；\n"
        "综合表现分析对所有作品进行整体评分；\n"
        "专业评价部分从商业、艺术和行业三个维度进行评估；\n"
        "发展建议为未来的艺术发展提供参考；\n"
        "最后是观众参与度分析，深入了解观众的互动行为。"
    )

def add_data_overview_slide(prs):
    """添加数据概况幻灯片"""
    slide_layout = prs.slide_layouts[1]  # 标题和内容布局
    slide = prs.slides.add_slide(slide_layout)
    
    # 设置背景
    set_slide_background(slide, "#0f3460")
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "数据概况"
    
    # 设置标题样式
    title_format = title.text_frame.paragraphs[0]
    title_format.font.color.rgb = RGBColor(255, 255, 255)  # 白色
    
    # 添加内容
    content.text = (
        "• 分析视频数量: 13首\n"
        "• 平均播放量: 2,019,763次\n"
        "• 平均弹幕数: 15,966条\n"
        "• 平均点赞数: 79,108个\n"
        "• 平均评论数: 9,107条"
    )
    
    # 设置内容样式
    for i, paragraph in enumerate(content.text_frame.paragraphs):
        paragraph.font.size = Pt(20)
        paragraph.font.color.rgb = RGBColor(220, 220, 220)  # 浅灰色
    
    # 添加备注
    slide.notes_slide.notes_text_frame.text = (
        "我们共分析了单依纯在《歌手》节目中的13首现场演唱视频。\n\n"
        "从基础数据来看，平均每首作品获得了超过200万次播放，近1.6万条弹幕，\n"
        "近8万个点赞和9千条评论，这表明单依纯的作品在B站平台具有很高的人气和关注度。\n\n"
        "这些数据为我们后续的深入分析提供了坚实的基础。"
    )

def add_heat_analysis_slide(prs):
    """添加热度分析幻灯片"""
    slide_layout = prs.slide_layouts[1]  # 标题和内容布局
    slide = prs.slides.add_slide(slide_layout)
    
    # 设置背景
    set_slide_background(slide, "#e94560")
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "热度趋势分析"
    
    # 设置标题样式
    title_format = title.text_frame.paragraphs[0]
    title_format.font.color.rgb = RGBColor(255, 255, 255)  # 白色
    
    # 添加内容
    content.text = (
        "• 最高播放量: 8,879,175次（《李白》）\n"
        "• 最低播放量: 467,563次（《爱情沙拉拉》）\n"
        "• 整体趋势: 后期作品热度持续上升\n"
        "• 观众互动: 弹幕、点赞、评论数量稳步增长"
    )
    
    # 设置内容样式
    for i, paragraph in enumerate(content.text_frame.paragraphs):
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = RGBColor(240, 240, 240)  # 浅灰色
    
    # 添加热度趋势分析图
    chart_path = "热度趋势分析_高级版.png"
    if os.path.exists(chart_path):
        left = Inches(6.5)
        top = Inches(1.5)
        width = Inches(6)
        height = Inches(4.5)
        slide.shapes.add_picture(chart_path, left, top, width, height)
    
    # 添加备注
    slide.notes_slide.notes_text_frame.text = (
        "热度分析显示，单依纯在《歌手》节目中的表现呈现稳步上升的趋势。\n\n"
        "其中《李白》以近890万次播放量成为最受欢迎的作品，而《爱情沙拉拉》播放量相对较低。\n\n"
        "整体来看，随着节目的推进，单依纯的作品热度持续上升，这可能与观众认知度提升、\n"
        "表演水平提高以及观众期待值增加等因素有关。\n\n"
        "观众互动数据也呈现稳步增长，表明观众对单依纯的表演越来越感兴趣。"
    )

def add_sentiment_analysis_slide(prs):
    """添加情感分析幻灯片"""
    slide_layout = prs.slide_layouts[1]  # 标题和内容布局
    slide = prs.slides.add_slide(slide_layout)
    
    # 设置背景
    set_slide_background(slide, "#25ccf7")
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "弹幕情感分析"
    
    # 设置标题样式
    title_format = title.text_frame.paragraphs[0]
    title_format.font.color.rgb = RGBColor(255, 255, 255)  # 白色
    
    # 添加内容
    content.text = (
        "• 正面评价: 占主导地位（70-80%）\n"
        "• 负面评价: 占比较少（1-5%）\n"
        "• 中性评价: 占比适中（15-25%）\n"
        "• 观众态度: 整体积极正面"
    )
    
    # 设置内容样式
    for i, paragraph in enumerate(content.text_frame.paragraphs):
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = RGBColor(240, 240, 240)  # 浅灰色
    
    # 添加弹幕情感分析图
    chart_path = "弹幕情感分析_高级版.png"
    if os.path.exists(chart_path):
        left = Inches(6.5)
        top = Inches(1.5)
        width = Inches(6)
        height = Inches(4.5)
        slide.shapes.add_picture(chart_path, left, top, width, height)
    
    # 添加备注
    slide.notes_slide.notes_text_frame.text = (
        "弹幕情感分析结果显示，观众对单依纯的表演整体持积极正面的态度。\n\n"
        "正面评价占比高达70-80%，主要包含'好听'、'棒'、'厉害'等赞美词汇，\n"
        "反映出观众对其演唱技巧和舞台表现的高度认可。\n\n"
        "负面评价仅占1-5%，主要集中在对音准或舞台表现细节的不同意见。\n"
        "中性评价占比15-25%，主要为对表演内容的客观描述。\n\n"
        "这种情感分布表明单依纯的表演获得了观众的广泛认可。"
    )

def add_keyword_analysis_slide(prs):
    """添加关键词分析幻灯片"""
    slide_layout = prs.slide_layouts[1]  # 标题和内容布局
    slide = prs.slides.add_slide(slide_layout)
    
    # 设置背景
    set_slide_background(slide, "#3b3b98")
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "关键词分析"
    
    # 设置标题样式
    title_format = title.text_frame.paragraphs[0]
    title_format.font.color.rgb = RGBColor(255, 255, 255)  # 白色
    
    # 添加内容
    content.text = (
        "高频正面词汇:\n"
        "• 好听、感谢、实至名归、回家、欢迎\n\n"
        "高频中性词汇:\n"
        "• 单依纯、成为、笑话、结局\n\n"
        "高频负面词汇:\n"
        "• 难听、差、烂"
    )
    
    # 设置内容样式
    for i, paragraph in enumerate(content.text_frame.paragraphs):
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = RGBColor(220, 220, 220)  # 浅灰色
    
    # 添加弹幕词云图
    chart_path = "弹幕词云_高级版.png"
    if os.path.exists(chart_path):
        left = Inches(6.5)
        top = Inches(1.5)
        width = Inches(6)
        height = Inches(4.5)
        slide.shapes.add_picture(chart_path, left, top, width, height)
    
    # 添加备注
    slide.notes_slide.notes_text_frame.text = (
        "通过词云分析，我们提取出观众讨论的高频词汇。\n\n"
        "正面词汇如'好听'、'感谢'、'实至名归'、'回家'、'欢迎'等，\n"
        "体现了观众对单依纯表演的高度认可和情感认同。\n\n"
        "中性词汇主要涉及表演者姓名和节目相关信息。\n\n"
        "负面词汇相对较少，主要为'难听'、'差'、'烂'等，\n"
        "表明观众整体满意度较高。\n\n"
        "这些关键词反映了观众关注的核心话题和评价维度。"
    )

def add_composite_score_slide(prs):
    """添加综合评分幻灯片"""
    slide_layout = prs.slide_layouts[1]  # 标题和内容布局
    slide = prs.slides.add_slide(slide_layout)
    
    # 设置背景
    set_slide_background(slide, "#f97f51")
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "综合表现分析"
    
    # 设置标题样式
    title_format = title.text_frame.paragraphs[0]
    title_format.font.color.rgb = RGBColor(255, 255, 255)  # 白色
    
    # 添加内容
    content.text = (
        "综合评分排名前五名:\n"
        "1. 《李白》 - 11,428,976分\n"
        "2. 《珠玉》歌手舞台 - 7,585,017分\n"
        "3. 《君》- 2,177,612分\n"
        "4. 《落叶归根》合作舞台 - 2,117,274分\n"
        "5. 《舞娘》- 1,972,426分"
    )
    
    # 设置内容样式
    for i, paragraph in enumerate(content.text_frame.paragraphs):
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = RGBColor(240, 240, 240)  # 浅灰色
    
    # 添加综合得分排名图
    chart_path = "视频综合得分排名_高级版.png"
    if os.path.exists(chart_path):
        left = Inches(6.5)
        top = Inches(1.5)
        width = Inches(6)
        height = Inches(4.5)
        slide.shapes.add_picture(chart_path, left, top, width, height)
    
    # 添加备注
    slide.notes_slide.notes_text_frame.text = (
        "综合表现分析通过多维度指标对视频进行评分，\n"
        "评分体系综合考虑了播放量、弹幕数、评论数、点赞数、投币数、收藏数和分享数等指标。\n\n"
        "《李白》以超过1142万分的高分位居榜首，表现最为突出。\n"
        "《珠玉》歌手舞台和《君》分别位列第二和第三名。\n\n"
        "前五名作品在各项指标上均表现优异，综合影响力最为突出。\n\n"
        "这种排名反映了观众对不同作品的偏好和认可程度。"
    )

def add_professional_evaluation_slide(prs):
    """添加专业评价幻灯片"""
    slide_layout = prs.slide_layouts[1]  # 标题和内容布局
    slide = prs.slides.add_slide(slide_layout)
    
    # 设置背景
    set_slide_background(slide, "#8b78e6")
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "专业评价"
    
    # 设置标题样式
    title_format = title.text_frame.paragraphs[0]
    title_format.font.color.rgb = RGBColor(255, 255, 255)  # 白色
    
    # 添加内容
    content.text = (
        "商业价值: ★★★★☆\n"
        "• 高关注度和传播效应\n"
        "• 具备商业合作潜力\n\n"
        "艺术价值: ★★★★☆\n"
        "• 出色的表演技巧和情感表达\n"
        "• 音乐风格多样性展现\n\n"
        "行业影响: ★★★★☆\n"
        "• 对节目和音乐市场有积极贡献"
    )
    
    # 设置内容样式
    for i, paragraph in enumerate(content.text_frame.paragraphs):
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = RGBColor(240, 240, 240)  # 浅灰色
    
    # 添加备注
    slide.notes_slide.notes_text_frame.text = (
        "基于以上数据分析结果，我们从商业价值、艺术价值和行业影响三个维度进行专业评价。\n\n"
        "商业价值方面，单依纯相关视频在B站平台获得大量关注，具有较高的商业开发潜力。\n\n"
        "艺术价值方面，观众普遍认为其演唱技巧娴熟，情感表达真挚动人，音乐风格多样。\n\n"
        "行业影响方面，其表演提升了节目的整体质量，对音乐类内容创作产生积极影响。\n\n"
        "总体而言，单依纯在《歌手》节目中的表现获得了专业认可。"
    )

def add_suggestions_slide(prs):
    """添加发展建议幻灯片"""
    slide_layout = prs.slide_layouts[1]  # 标题和内容布局
    slide = prs.slides.add_slide(slide_layout)
    
    # 设置背景
    set_slide_background(slide, "#00d2d3")
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "发展建议"
    
    # 设置标题样式
    title_format = title.text_frame.paragraphs[0]
    title_format.font.color.rgb = RGBColor(255, 255, 255)  # 白色
    
    # 添加内容
    content.text = (
        "1. 保持高质量的现场表演\n"
        "2. 尝试更多元化的音乐风格\n"
        "3. 加强与观众的互动\n"
        "4. 利用高人气进行商业合作\n"
        "5. 注重作品创新"
    )
    
    # 设置内容样式
    for i, paragraph in enumerate(content.text_frame.paragraphs):
        paragraph.font.size = Pt(20)
        paragraph.font.color.rgb = RGBColor(240, 240, 240)  # 浅灰色
    
    # 添加备注
    slide.notes_slide.notes_text_frame.text = (
        "基于数据分析结果和专业评价，我们为单依纯未来的艺术发展提出以下建议：\n\n"
        "首先，继续保持高质量的现场表演，巩固专业实力和观众口碑；\n"
        "其次，尝试更多元化的音乐风格，拓展艺术表现领域；\n"
        "第三，加强与观众的互动，提升粉丝粘性和社区活跃度；\n"
        "第四，利用高人气进行商业合作和品牌推广，实现艺术与商业的良性循环；\n"
        "最后，在保持质量的基础上，注重作品创新。\n\n"
        "这些建议旨在帮助单依纯在艺术道路上取得更大的成就。"
    )

def add_engagement_analysis_slide(prs):
    """添加观众参与度分析幻灯片"""
    slide_layout = prs.slide_layouts[1]  # 标题和内容布局
    slide = prs.slides.add_slide(slide_layout)
    
    # 设置背景
    set_slide_background(slide, "#ff9ff3")
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "观众参与度分析"
    
    # 设置标题样式
    title_format = title.text_frame.paragraphs[0]
    title_format.font.color.rgb = RGBColor(255, 255, 255)  # 白色
    
    # 添加内容
    content.text = (
        "参与度计算公式:\n"
        "参与度 = (弹幕数 + 评论数×2 + 点赞数×0.5 + 投币数×3 + 收藏数×3 + 分享数×4) / 播放量 × 100%\n\n"
        "排名情况:\n"
        "1. 《君》 - 26.46%\n"
        "2. 《有趣》 - 25.62%\n"
        "3. 《李白》 - 14.57%"
    )
    
    # 设置内容样式
    for i, paragraph in enumerate(content.text_frame.paragraphs):
        paragraph.font.size = Pt(14)
        paragraph.font.color.rgb = RGBColor(240, 240, 240)  # 浅灰色
    
    # 添加观众参与度分析图
    chart_path = "观众参与度分析.png"
    if os.path.exists(chart_path):
        left = Inches(6.5)
        top = Inches(1.5)
        width = Inches(6)
        height = Inches(4.5)
        slide.shapes.add_picture(chart_path, left, top, width, height)
    
    # 添加备注
    slide.notes_slide.notes_text_frame.text = (
        "观众参与度是衡量内容吸引力和观众粘性的重要指标。\n\n"
        "我们通过综合考虑弹幕、评论、点赞、投币、收藏和分享等多种互动行为，\n"
        "并将其与播放量进行标准化处理，计算出观众参与度。\n\n"
        "《君》以26.46%的参与度位居榜首，表明该作品不仅获得了高播放量，\n"
        "更重要的是激发了观众的深度互动。\n\n"
        "《有趣》和《李白》分别位列第二和第三名。\n\n"
        "这一分析为内容创作者提供了有价值的参考，帮助他们了解如何创作更能激发观众互动的内容。"
    )

def add_conclusion_slide(prs):
    """添加结论幻灯片"""
    slide_layout = prs.slide_layouts[1]  # 标题和内容布局
    slide = prs.slides.add_slide(slide_layout)
    
    # 设置背景
    set_slide_background(slide, "#54a0ff")
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "主要结论"
    
    # 设置标题样式
    title_format = title.text_frame.paragraphs[0]
    title_format.font.color.rgb = RGBColor(255, 255, 255)  # 白色
    
    # 添加内容
    content.text = (
        "1. 单依纯在《歌手》节目中表现出色，热度持续上升\n"
        "2. 观众对其表演普遍给予高度评价，正面反馈占主导\n"
        "3. 不同音乐风格的作品均获得较好反响，展现较强适应性\n"
        "4. 互动数据表现优异，粉丝粘性和传播效应明显\n"
        "5. 专业评价获得认可，具备良好的商业和艺术价值"
    )
    
    # 设置内容样式
    for i, paragraph in enumerate(content.text_frame.paragraphs):
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = RGBColor(240, 240, 240)  # 浅灰色
    
    # 添加备注
    slide.notes_slide.notes_text_frame.text = (
        "通过全面的数据分析，我们得出以下主要结论：\n\n"
        "单依纯在《歌手》节目中的表现非常出色，作品热度呈现持续上升趋势；\n"
        "观众对其表演普遍给予高度评价，正面反馈占据主导地位；\n"
        "她在不同音乐风格的作品中均有良好表现，展现了较强的适应性；\n"
        "互动数据表现优异，显示出良好的粉丝粘性和传播效应；\n"
        "专业评价获得认可，具备良好的商业和艺术价值。\n\n"
        "这些结论为单依纯未来的发展提供了有力的数据支持。"
    )

def add_final_slide(prs):
    """添加结束幻灯片"""
    slide_layout = prs.slide_layouts[0]  # 标题幻灯片布局
    slide = prs.slides.add_slide(slide_layout)
    
    # 设置背景
    set_slide_background(slide, "#5f27cd")
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "谢谢观看"
    subtitle.text = "单依纯《歌手》节目数据分析报告"
    
    # 设置标题样式
    title_format = title.text_frame.paragraphs[0]
    title_format.font.size = Pt(44)
    title_format.font.bold = True
    title_format.font.color.rgb = RGBColor(255, 255, 255)  # 白色
    
    subtitle_format = subtitle.text_frame.paragraphs[0]
    subtitle_format.font.size = Pt(28)
    subtitle_format.font.color.rgb = RGBColor(220, 220, 220)  # 浅灰色
    
    # 添加装饰性元素
    left = Inches(1)
    top = Inches(0.5)
    width = Inches(2)
    height = Inches(2)
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, width, height)
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(220, 20, 60)  # 猩红色
    shape.line.color.rgb = RGBColor(255, 255, 255)  # 白色边框
    
    # 添加备注
    slide.notes_slide.notes_text_frame.text = (
        "感谢您观看本次数据分析报告演示。\n\n"
        "如果您对报告内容有任何疑问或需要进一步了解某些分析细节，\n"
        "欢迎随时提出问题。\n\n"
        "希望这份报告能为单依纯未来的艺术发展提供有价值的参考。"
    )

def main():
    """主函数"""
    print("开始将PDF报告转换为PPT格式...")
    
    # 创建演示文稿
    prs = create_presentation()
    
    # 添加各幻灯片
    add_title_slide(prs)
    add_outline_slide(prs)
    add_data_overview_slide(prs)
    add_heat_analysis_slide(prs)
    add_sentiment_analysis_slide(prs)
    add_keyword_analysis_slide(prs)
    add_composite_score_slide(prs)
    add_professional_evaluation_slide(prs)
    add_suggestions_slide(prs)
    add_engagement_analysis_slide(prs)
    add_conclusion_slide(prs)
    add_final_slide(prs)
    
    # 保存演示文稿
    output_file = "单依纯_歌手节目数据分析报告_带图表.pptx"
    prs.save(output_file)
    print(f"PPT报告已生成: {output_file}")

if __name__ == "__main__":
    main()